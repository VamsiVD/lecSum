import json
import os
import urllib.parse
import uuid
from datetime import datetime, timezone

import boto3

s3_client = boto3.client("s3")
transcribe_client = boto3.client("transcribe")
textract_client = boto3.client("textract")
dynamodb = boto3.resource("dynamodb")

TABLE_NAME = os.environ.get("DYNAMODB_TABLE", "lecsum-jobs")
OUTPUT_BUCKET = os.environ.get("OUTPUT_BUCKET", "lecsum-transcripts-text")
TRANSCRIBE_OUTPUT_BUCKET = os.environ.get("TRANSCRIBE_OUTPUT_BUCKET", "lecsum-transcripts-outbox")

AUDIO_FORMATS = {"mp3", "wav", "flac", "ogg", "amr", "webm"}
AUDIO_FORMAT_MAP = {
    "mp3": "mp3",
    "wav": "wav",
    "flac": "flac",
    "ogg": "ogg",
    "amr": "amr",
    "webm": "webm",
    "m4a": "mp4",
    "mp4": "mp4",
}
PDF_FORMATS = {"pdf"}
WORD_FORMATS = {"docx"}
PPTX_FORMATS = {"pptx"}
IMAGE_FORMATS = {"jpg", "jpeg", "png", "tiff", "tif"}


def write_dynamo(upload_key: str, file_name: str, status: str, job_name: str = "", extra: dict = None):
    table = dynamodb.Table(TABLE_NAME)
    item = {
        "uploadKey": upload_key,
        "fileName": file_name,
        "status": status,
        "createdAt": datetime.now(timezone.utc).isoformat(),
    }
    if job_name:
        item["jobName"] = job_name
    if extra:
        item.update(extra)
    table.put_item(Item=item)


def update_dynamo(upload_key: str, status: str, transcript_key: str = ""):
    table = dynamodb.Table(TABLE_NAME)
    expr = "SET #s = :s"
    names = {"#s": "status"}
    values = {":s": status}
    if transcript_key:
        expr += ", transcriptKey = :tk"
        values[":tk"] = transcript_key
    table.update_item(
        Key={"uploadKey": upload_key},
        UpdateExpression=expr,
        ExpressionAttributeNames=names,
        ExpressionAttributeValues=values,
    )


def save_text(text: str, upload_key: str) -> str:
    """Save extracted text to the output bucket, return the S3 key."""
    base = upload_key.rsplit(".", 1)[0]
    out_key = f"{base}.txt"
    s3_client.put_object(
        Bucket=OUTPUT_BUCKET,
        Key=out_key,
        Body=text.encode("utf-8"),
        ContentType="text/plain",
    )
    return out_key


def handle_audio(bucket: str, key: str, upload_key: str, file_name: str, ext: str):
    media_format = AUDIO_FORMAT_MAP.get(ext, "mp3")
    job_name = f"lecsum-job-{uuid.uuid4().hex[:8]}"
    file_uri = f"s3://{bucket}/{key}"

    write_dynamo(upload_key, file_name, "transcribing", job_name)

    transcribe_client.start_transcription_job(
        TranscriptionJobName=job_name,
        Media={"MediaFileUri": file_uri},
        MediaFormat=media_format,
        LanguageCode="en-US",
        OutputBucketName=TRANSCRIBE_OUTPUT_BUCKET,
    )
    print(f"Audio: started Transcribe job {job_name}")


def handle_pdf(bucket: str, key: str, upload_key: str, file_name: str):
    import io

    try:
        import pdfplumber
    except ImportError:
        import pypdf as _pypdf  # fallback

        obj = s3_client.get_object(Bucket=bucket, Key=key)
        reader = _pypdf.PdfReader(io.BytesIO(obj["Body"].read()))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
    else:
        obj = s3_client.get_object(Bucket=bucket, Key=key)
        with pdfplumber.open(io.BytesIO(obj["Body"].read())) as pdf:
            text = "\n".join(page.extract_text() or "" for page in pdf.pages)

    write_dynamo(upload_key, file_name, "extracting")
    out_key = save_text(text.strip(), upload_key)
    update_dynamo(upload_key, "done", out_key)
    print(f"PDF: extracted {len(text)} chars → {out_key}")


def handle_docx(bucket: str, key: str, upload_key: str, file_name: str):
    import io

    from docx import Document

    obj = s3_client.get_object(Bucket=bucket, Key=key)
    doc = Document(io.BytesIO(obj["Body"].read()))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    write_dynamo(upload_key, file_name, "extracting")
    out_key = save_text(text.strip(), upload_key)
    update_dynamo(upload_key, "done", out_key)
    print(f"DOCX: extracted {len(text)} chars → {out_key}")


def handle_pptx(bucket: str, key: str, upload_key: str, file_name: str):
    import io

    from pptx import Presentation

    obj = s3_client.get_object(Bucket=bucket, Key=key)
    prs = Presentation(io.BytesIO(obj["Body"].read()))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))
    text = "\n\n".join(slides)

    write_dynamo(upload_key, file_name, "extracting")
    out_key = save_text(text.strip(), upload_key)
    update_dynamo(upload_key, "done", out_key)
    print(f"PPTX: extracted {len(slides)} slides → {out_key}")


def handle_image(bucket: str, key: str, upload_key: str, file_name: str):
    write_dynamo(upload_key, file_name, "extracting")

    response = textract_client.detect_document_text(
        Document={"S3Object": {"Bucket": bucket, "Name": key}}
    )
    lines = [
        block["Text"]
        for block in response["Blocks"]
        if block["BlockType"] == "LINE"
    ]
    text = "\n".join(lines)

    out_key = save_text(text.strip(), upload_key)
    update_dynamo(upload_key, "done", out_key)
    print(f"Image: extracted {len(lines)} lines via Textract → {out_key}")


def lambda_handler(event, context):
    try:
        bucket = event["Records"][0]["s3"]["bucket"]["name"]
        key = urllib.parse.unquote_plus(
            event["Records"][0]["s3"]["object"]["key"], encoding="utf-8"
        )
    except Exception as e:
        print(f"Error parsing S3 event: {e}")
        return {"statusCode": 400, "body": "Invalid S3 event"}

    ext = key.rsplit(".", 1)[-1].lower() if "." in key else ""
    file_name = key.split("/")[-1]
    upload_key = key

    print(f"Router received: {key} (ext={ext})")

    try:
        if ext in AUDIO_FORMAT_MAP:
            handle_audio(bucket, key, upload_key, file_name, ext)
        elif ext in PDF_FORMATS:
            handle_pdf(bucket, key, upload_key, file_name)
        elif ext in WORD_FORMATS:
            handle_docx(bucket, key, upload_key, file_name)
        elif ext in PPTX_FORMATS:
            handle_pptx(bucket, key, upload_key, file_name)
        elif ext in IMAGE_FORMATS:
            handle_image(bucket, key, upload_key, file_name)
        else:
            print(f"Unsupported format: {ext}")
            return {"statusCode": 400, "body": f"Unsupported format: {ext}"}

        return {"statusCode": 200, "body": f"Processed {file_name}"}

    except Exception as e:
        print(f"Error processing {file_name}: {e}")
        update_dynamo(upload_key, "error")
        return {"statusCode": 500, "body": str(e)}